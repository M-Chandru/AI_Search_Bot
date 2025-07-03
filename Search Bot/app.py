#'app.py - A Flask application for file management and AI-based question answering using Google Gemini API.'
import os, fitz, docx2txt, pandas as pd, traceback, json
from pptx import Presentation
from flask import Flask, request, jsonify, render_template, redirect, url_for, session
from flask_cors import CORS
from dotenv import load_dotenv
import google.generativeai as genai

# Initialize Flask app
app = Flask(__name__)
app.secret_key = 'your_secret_key'
CORS(app)

# Folder and config setup
UPLOAD_FOLDER = 'uploads'
USER_FILE = 'users.json'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# Load environment variables and Gemini API key
load_dotenv()
API_KEY = os.getenv("API_KEY")
genai.configure(api_key=API_KEY)
model = genai.GenerativeModel("gemini-1.5-flash")

# Allowed file types
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'xlsx', 'pptx'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# User handling functions
def load_users():
    if not os.path.exists(USER_FILE):
        return {}
    with open(USER_FILE, 'r') as f:
        return json.load(f)

def save_users(users):
    with open(USER_FILE, 'w') as f:
        json.dump(users, f)

# File parsing functions
def parse_pdf(path):
    try:
        doc = fitz.open(path)
        return "\n".join([page.get_text() for page in doc])
    except Exception as e:
        return f"[Error parsing PDF: {e}]"

def parse_docx(path):
    try:
        return docx2txt.process(path)
    except Exception as e:
        return f"[Error parsing DOCX: {e}]"

def parse_xlsx(path):
    try:
        df = pd.read_excel(path)
        return df.to_string()
    except Exception as e:
        return f"[Error parsing XLSX: {e}]"

def parse_pptx(path):
    try:
        text = ""
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"[Error parsing PPTX: {e}]"

# Routes
@app.route('/')
def index():
    return render_template("index.html")

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        users = load_users()

        if username in users:
            return render_template("register.html", error="Username already exists")

        users[username] = password
        save_users(users)
        return redirect(url_for('index'))
    return render_template("register.html")

@app.route('/login', methods=['POST'])
def login():
    username = request.form['username']
    password = request.form['password']
    users = load_users()

    if username in users and users[username] == password:
        session['username'] = username
        session['role'] = 'admin' if username == 'admin' else 'user'
        return redirect(url_for('admin' if username == 'admin' else 'user'))
    return render_template("index.html", error="Invalid credentials")

@app.route('/admin')
def admin():
    if session.get('role') == 'admin':
        return render_template('admin.html', role='admin')
    return redirect(url_for('index'))

@app.route('/user')
def user():
    if session.get('role') == 'user':
        return render_template('user.html', role='user')
    return redirect(url_for('index'))

@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('index'))

@app.route('/upload', methods=['POST'])
def upload_file():
    if session.get('role') != 'admin':
        return jsonify({'error': 'Unauthorized'}), 403
    file = request.files.get('file')
    if file and allowed_file(file.filename):
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(filepath)
        return jsonify({'message': 'File uploaded successfully'})
    return jsonify({'error': 'Invalid file'})

@app.route('/delete', methods=['POST'])
def delete_file():
    if session.get('role') != 'admin':
        return jsonify({'error': 'Unauthorized'}), 403
    filename = request.form.get('filename')
    path = os.path.join(UPLOAD_FOLDER, filename)
    if os.path.exists(path):
        os.remove(path)
        return jsonify({'message': 'File deleted'})
    return jsonify({'error': 'File not found'})

@app.route('/list', methods=['GET'])
def list_files():
    return jsonify({'files': os.listdir(UPLOAD_FOLDER)})

@app.route('/ask', methods=['POST'])
def ask():
    question = request.form.get('question')
    context = ""
    for file in os.listdir(UPLOAD_FOLDER):
        path = os.path.join(UPLOAD_FOLDER, file)
        try:
            if file.endswith(".pdf"):
                context += parse_pdf(path)
            elif file.endswith(".docx"):
                context += parse_docx(path)
            elif file.endswith(".xlsx"):
                context += parse_xlsx(path)
            elif file.endswith(".pptx"):
                context += parse_pptx(path)
        except Exception as e:
            context += f"[Error reading {file}: {e}]\n"
    prompt = f"Use the following document context to answer the question.\n\nContext:\n{context[:8000]}\n\nQuestion: {question}"
    try:
        response = model.generate_content(prompt)
        answer = response.text if hasattr(response, 'text') else "No response."
        return jsonify({'answer': answer})
    except Exception as e:
        traceback.print_exc()
        return jsonify({'error': str(e)})

if __name__ == '__main__':
    app.run(debug=True, port=8080)
